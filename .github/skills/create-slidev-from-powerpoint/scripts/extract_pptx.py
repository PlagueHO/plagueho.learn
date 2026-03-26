"""
Extract theme, slides, and images from a PowerPoint file for Slidev conversion.

Usage:
    python extract_pptx.py <input.pptx> <output_dir>

Outputs:
    <output_dir>/theme.json     - Colors, fonts, backgrounds
    <output_dir>/slides.json    - Slide content, layouts, notes
    <output_dir>/images/        - Extracted media files
    <output_dir>/layouts.json   - Slide layout definitions
"""

import json
import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def emu_to_px(emu):
    """Convert EMU (English Metric Units) to approximate pixels at 96 DPI."""
    if emu is None:
        return None
    return round(emu / 914400 * 96)


def emu_to_inches(emu):
    """Convert EMU to inches."""
    if emu is None:
        return None
    return round(emu / 914400, 2)


def rgb_to_hex(rgb_color):
    """Convert an RGBColor or tuple to hex string."""
    if rgb_color is None:
        return None
    if isinstance(rgb_color, RGBColor):
        return f"#{rgb_color}"
    if isinstance(rgb_color, tuple) and len(rgb_color) == 3:
        return f"#{rgb_color[0]:02x}{rgb_color[1]:02x}{rgb_color[2]:02x}"
    return str(rgb_color)


def extract_theme_colors_from_xml(prs):
    """Extract theme color scheme from the presentation XML."""
    colors = {}
    try:
        from lxml import etree
    except ImportError:
        import xml.etree.ElementTree as etree

    theme_elem = prs.slide_masters[0].element
    nsmap = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    }

    # Get theme element from slide master
    slide_master = prs.slide_masters[0]
    theme_part = slide_master.part.slide_master.element

    # Try to access the theme XML directly
    for rel in slide_master.part.rels.values():
        if "theme" in rel.reltype:
            theme_xml = rel.target_part.blob
            root = etree.fromstring(theme_xml) if hasattr(etree, 'fromstring') else etree.XML(theme_xml)

            # Define namespace
            ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

            # Extract color scheme
            clrScheme = root.find('.//a:clrScheme', ns)
            if clrScheme is not None:
                colors['scheme_name'] = clrScheme.get('name', 'Unknown')
                for child in clrScheme:
                    tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                    srgb = child.find('a:srgbClr', ns)
                    sys_clr = child.find('a:sysClr', ns)
                    if srgb is not None:
                        colors[tag] = f"#{srgb.get('val')}"
                    elif sys_clr is not None:
                        colors[tag] = f"#{sys_clr.get('lastClr', sys_clr.get('val', '000000'))}"

            # Extract font scheme
            fontScheme = root.find('.//a:fontScheme', ns)
            if fontScheme is not None:
                colors['font_scheme_name'] = fontScheme.get('name', 'Unknown')
                major_font = fontScheme.find('.//a:majorFont/a:latin', ns)
                minor_font = fontScheme.find('.//a:minorFont/a:latin', ns)
                if major_font is not None:
                    colors['major_font'] = major_font.get('typeface')
                if minor_font is not None:
                    colors['minor_font'] = minor_font.get('typeface')

                # Also get East Asian and Complex Script fonts
                major_ea = fontScheme.find('.//a:majorFont/a:ea', ns)
                minor_ea = fontScheme.find('.//a:minorFont/a:ea', ns)
                if major_ea is not None:
                    colors['major_font_ea'] = major_ea.get('typeface')
                if minor_ea is not None:
                    colors['minor_font_ea'] = minor_ea.get('typeface')

            # Extract format scheme / background fill styles
            fmtScheme = root.find('.//a:fmtScheme', ns)
            if fmtScheme is not None:
                colors['format_scheme_name'] = fmtScheme.get('name', 'Unknown')

            break

    return colors


def extract_slide_background(slide):
    """Extract the background info from a slide."""
    bg_info = {}
    try:
        bg = slide.background
        fill = bg.fill
        if fill.type is not None:
            bg_info['fill_type'] = str(fill.type)
            if fill.type == 1:  # Solid
                try:
                    bg_info['color'] = rgb_to_hex(fill.fore_color.rgb)
                except Exception:
                    bg_info['color'] = None
    except Exception:
        pass
    return bg_info


def extract_text_formatting(run):
    """Extract text formatting from a run."""
    fmt = {}
    try:
        font = run.font
        if font.bold:
            fmt['bold'] = True
        if font.italic:
            fmt['italic'] = True
        if font.size is not None:
            fmt['size_pt'] = font.size.pt
        if font.color and font.color.rgb:
            fmt['color'] = rgb_to_hex(font.color.rgb)
        if font.name:
            fmt['font'] = font.name
        if font.underline:
            fmt['underline'] = True
    except Exception:
        pass
    return fmt


def extract_paragraph(paragraph):
    """Extract a paragraph with its runs and formatting."""
    para_info = {
        'text': paragraph.text,
        'runs': [],
        'level': paragraph.level if paragraph.level else 0,
    }

    # Alignment
    if paragraph.alignment is not None:
        align_map = {
            PP_ALIGN.LEFT: 'left',
            PP_ALIGN.CENTER: 'center',
            PP_ALIGN.RIGHT: 'right',
            PP_ALIGN.JUSTIFY: 'justify',
        }
        para_info['alignment'] = align_map.get(paragraph.alignment, str(paragraph.alignment))

    for run in paragraph.runs:
        run_info = {
            'text': run.text,
            'formatting': extract_text_formatting(run),
        }
        para_info['runs'].append(run_info)

    return para_info


def extract_shape(shape):
    """Extract info from a shape."""
    shape_info = {
        'shape_type': str(shape.shape_type) if shape.shape_type else 'unknown',
        'name': shape.name,
        'left': emu_to_px(shape.left),
        'top': emu_to_px(shape.top),
        'width': emu_to_px(shape.width),
        'height': emu_to_px(shape.height),
    }

    # Check if it's a placeholder
    if shape.is_placeholder:
        ph = shape.placeholder_format
        shape_info['is_placeholder'] = True
        shape_info['placeholder_type'] = str(ph.type) if ph.type else None
        shape_info['placeholder_idx'] = ph.idx

    # Extract text
    if shape.has_text_frame:
        shape_info['has_text'] = True
        shape_info['paragraphs'] = []
        for para in shape.text_frame.paragraphs:
            shape_info['paragraphs'].append(extract_paragraph(para))
        shape_info['full_text'] = shape.text_frame.text

    # Extract image
    if shape.shape_type == 13:  # Picture
        try:
            image = shape.image
            shape_info['has_image'] = True
            shape_info['image_content_type'] = image.content_type
            ext = image.content_type.split('/')[-1]
            if ext == 'jpeg':
                ext = 'jpg'
            elif ext == 'svg+xml':
                ext = 'svg'
            shape_info['image_ext'] = ext
        except Exception:
            pass

    # Extract table
    if shape.has_table:
        table = shape.table
        shape_info['has_table'] = True
        shape_info['table_rows'] = []
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                row_data.append(cell.text)
            shape_info['table_rows'].append(row_data)

    return shape_info


def extract_notes(slide):
    """Extract speaker notes from a slide."""
    try:
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            return notes_slide.notes_text_frame.text
    except Exception:
        pass
    return ""


def extract_slide(slide, slide_num):
    """Extract all content from a slide."""
    slide_info = {
        'number': slide_num,
        'layout_name': slide.slide_layout.name if slide.slide_layout else 'Unknown',
        'background': extract_slide_background(slide),
        'shapes': [],
        'notes': extract_notes(slide),
    }

    for shape in slide.shapes:
        slide_info['shapes'].append(extract_shape(shape))

    return slide_info


def extract_images(prs, output_dir):
    """Extract all images from the presentation."""
    images_dir = os.path.join(output_dir, 'images')
    os.makedirs(images_dir, exist_ok=True)

    image_map = {}
    img_counter = 0

    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture
                try:
                    image = shape.image
                    ext = image.content_type.split('/')[-1]
                    if ext == 'jpeg':
                        ext = 'jpg'
                    elif ext == 'svg+xml':
                        ext = 'svg'
                    img_counter += 1
                    filename = f"slide{slide_num}_img{img_counter}.{ext}"
                    filepath = os.path.join(images_dir, filename)
                    with open(filepath, 'wb') as f:
                        f.write(image.blob)
                    image_map[shape.name] = filename
                    print(f"  Extracted: {filename} ({shape.width}x{shape.height} EMU)")
                except Exception as e:
                    print(f"  Warning: Could not extract image from {shape.name}: {e}")

    return image_map


def extract_slide_layouts(prs):
    """Extract available slide layout definitions."""
    layouts = []
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            layout_info = {
                'name': layout.name,
                'placeholders': [],
            }
            for ph in layout.placeholders:
                layout_info['placeholders'].append({
                    'idx': ph.placeholder_format.idx,
                    'type': str(ph.placeholder_format.type),
                    'name': ph.name,
                    'left': emu_to_px(ph.left),
                    'top': emu_to_px(ph.top),
                    'width': emu_to_px(ph.width),
                    'height': emu_to_px(ph.height),
                })
            layouts.append(layout_info)
    return layouts


def extract_presentation_info(prs):
    """Extract overall presentation info."""
    info = {
        'slide_width_px': emu_to_px(prs.slide_width),
        'slide_height_px': emu_to_px(prs.slide_height),
        'slide_width_inches': emu_to_inches(prs.slide_width),
        'slide_height_inches': emu_to_inches(prs.slide_height),
        'slide_count': len(prs.slides),
    }
    return info


def main():
    if len(sys.argv) < 3:
        print(f"Usage: python {sys.argv[0]} <input.pptx> <output_dir>")
        sys.exit(1)

    pptx_path = sys.argv[1]
    output_dir = sys.argv[2]

    if not os.path.exists(pptx_path):
        print(f"Error: File not found: {pptx_path}")
        sys.exit(1)

    os.makedirs(output_dir, exist_ok=True)

    print(f"Loading: {pptx_path}")
    prs = Presentation(pptx_path)

    # Presentation info
    print("Extracting presentation info...")
    pres_info = extract_presentation_info(prs)
    print(f"  Size: {pres_info['slide_width_px']}x{pres_info['slide_height_px']}px "
          f"({pres_info['slide_width_inches']}x{pres_info['slide_height_inches']}in)")
    print(f"  Slides: {pres_info['slide_count']}")

    # Theme
    print("Extracting theme...")
    theme = extract_theme_colors_from_xml(prs)
    theme.update(pres_info)
    with open(os.path.join(output_dir, 'theme.json'), 'w', encoding='utf-8') as f:
        json.dump(theme, f, indent=2, ensure_ascii=False)
    print(f"  Colors: {', '.join(k + '=' + v for k, v in theme.items() if k.startswith('dk') or k.startswith('lt') or k.startswith('accent'))}")
    print(f"  Fonts: major={theme.get('major_font')}, minor={theme.get('minor_font')}")

    # Layouts
    print("Extracting layouts...")
    layouts = extract_slide_layouts(prs)
    with open(os.path.join(output_dir, 'layouts.json'), 'w', encoding='utf-8') as f:
        json.dump(layouts, f, indent=2, ensure_ascii=False)
    print(f"  Found {len(layouts)} layouts: {', '.join(l['name'] for l in layouts[:10])}...")

    # Images
    print("Extracting images...")
    image_map = extract_images(prs, output_dir)
    print(f"  Extracted {len(image_map)} images")

    # Slides
    print("Extracting slides...")
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        slide_info = extract_slide(slide, i)
        slides.append(slide_info)
        text_preview = ""
        for shape in slide_info['shapes']:
            if 'full_text' in shape and shape['full_text']:
                text_preview = shape['full_text'][:60]
                break
        print(f"  Slide {i}: layout='{slide_info['layout_name']}' "
              f"shapes={len(slide_info['shapes'])} "
              f"text='{text_preview}...'")

    # Add image filenames to slide shapes
    for slide_data in slides:
        for shape in slide_data['shapes']:
            if shape['name'] in image_map:
                shape['image_filename'] = image_map[shape['name']]

    with open(os.path.join(output_dir, 'slides.json'), 'w', encoding='utf-8') as f:
        json.dump(slides, f, indent=2, ensure_ascii=False)

    print(f"\nDone! Output written to: {output_dir}")
    print(f"  theme.json   - Theme colors and fonts")
    print(f"  layouts.json - {len(layouts)} slide layouts")
    print(f"  slides.json  - {len(slides)} slides with content")
    print(f"  images/      - {len(image_map)} media files")


if __name__ == '__main__':
    main()
